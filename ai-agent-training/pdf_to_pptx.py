# -*- coding: utf-8 -*-
"""把 exports/*.pdf 每页渲染为高清图并封装成 16:9 PPTX（用于现场放映）"""
import os, io, glob
import fitz
from pptx import Presentation
from pptx.util import Emu

ROOT = r"C:\Users\think\WorkBuddy\aippt8\integrated"
PDF_DIR = os.path.join(ROOT, "exports")
OUT_DIR = os.path.join(ROOT, "pptx")
os.makedirs(OUT_DIR, exist_ok=True)

W = Emu(12192000)   # 13.333in = 16:9 宽屏
H = Emu(6858000)    # 7.5in
ZOOM = 2.0          # 960*2 = 1920px 宽，够投影

for pdf_path in sorted(glob.glob(os.path.join(PDF_DIR, "*.pdf"))):
    name = os.path.splitext(os.path.basename(pdf_path))[0]
    doc = fitz.open(pdf_path)
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    blank = prs.slide_layouts[6]
    for page in doc:
        pix = page.get_pixmap(matrix=fitz.Matrix(ZOOM, ZOOM), alpha=False)
        img = io.BytesIO(pix.tobytes("png"))
        s = prs.slides.add_slide(blank)
        s.shapes.add_picture(img, 0, 0, width=W, height=H)
    out = os.path.join(OUT_DIR, name + "（放映版）.pptx")
    prs.save(out)
    print("OK  %-32s %2d 页  %.1f MB" % (name, len(doc), os.path.getsize(out) / 1048576))
    doc.close()
